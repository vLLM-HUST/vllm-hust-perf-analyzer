#!/usr/bin/env python3
"""Generate TraceLoom introduction PPTX slide with dark theme."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# --- Colors (dark theme) ---
BG_DARK = RGBColor(0x0D, 0x11, 0x17)
BG_CARD = RGBColor(0x16, 0x1B, 0x22)
BORDER = RGBColor(0x30, 0x36, 0x3D)
TEXT_PRIMARY = RGBColor(0xF0, 0xF6, 0xFC)
TEXT_SECONDARY = RGBColor(0xC9, 0xD1, 0xD9)
TEXT_MUTED = RGBColor(0x8B, 0x94, 0x9E)
BLUE = RGBColor(0x58, 0xA6, 0xFF)
GREEN = RGBColor(0x7E, 0xE7, 0x87)
PURPLE = RGBColor(0xD2, 0xA8, 0xFF)
GRADIENT_TOP = RGBColor(0x1A, 0x1F, 0x2E)


def set_slide_bg(slide, color):
    """Set slide background to solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=12,
                 bold=False, color=TEXT_PRIMARY, alignment=PP_ALIGN.LEFT,
                 font_name="Arial"):
    """Add a simple text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rich_paragraph(tf, runs, space_before=0, space_after=0, alignment=PP_ALIGN.LEFT):
    """Add a paragraph with multiple formatted runs."""
    p = tf.add_paragraph()
    p.alignment = alignment
    if space_before:
        p.space_before = Pt(space_before)
    if space_after:
        p.space_after = Pt(space_after)
    for run_text, run_props in runs:
        run = p.add_run()
        run.text = run_text
        run.font.size = Pt(run_props.get("size", 12))
        run.font.bold = run_props.get("bold", False)
        run.font.color.rgb = run_props.get("color", TEXT_PRIMARY)
        run.font.name = run_props.get("font", "Arial")
    return p


def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # Reduce corner radius
    shape.adjustments[0] = 0.05
    return shape


def add_gradient_bar(slide, left, top, width, height):
    """Add the top gradient accent bar."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()
    return shape


def build_slide(prs):
    """Build the single TraceLoom introduction slide."""
    slide_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide, BG_DARK)

    # Dimensions
    sw = prs.slide_width  # 13.333 inches (widescreen)
    sh = prs.slide_height  # 7.5 inches

    # === TOP GRADIENT BAR ===
    add_gradient_bar(slide, Inches(0), Inches(0), sw, Pt(5))

    # === HEADER ===
    # Title: TraceLoom
    add_text_box(slide, Inches(0.7), Inches(0.4), Inches(3), Inches(0.7),
                 "TraceLoom", font_size=36, bold=True, color=TEXT_PRIMARY)

    # Subtitle
    add_text_box(slide, Inches(3.6), Inches(0.55), Inches(4), Inches(0.5),
                 "离线加速器性能分析工具", font_size=16, color=TEXT_MUTED)

    # Badge - MIT License
    badge = add_rounded_rect(slide, Inches(10.8), Inches(0.5), Inches(1.2), Inches(0.35),
                             RGBColor(0x12, 0x24, 0x3A), BLUE)
    badge.text_frame.paragraphs[0].text = "MIT License"
    badge.text_frame.paragraphs[0].font.size = Pt(9)
    badge.text_frame.paragraphs[0].font.color.rgb = BLUE
    badge.text_frame.paragraphs[0].font.bold = True
    badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    badge.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # === SEPARATOR LINE ===
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.1), Inches(11.9), Pt(1)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BORDER
    line.line.fill.background()

    # ===================== LEFT COLUMN =====================
    left_x = Inches(0.7)

    # --- Tagline / Positioning ---
    # Blue left border accent
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left_x, Inches(1.35), Pt(4), Inches(0.65)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = BLUE
    accent_bar.line.fill.background()

    tagline_box = slide.shapes.add_textbox(Inches(0.95), Inches(1.35), Inches(5.3), Inches(0.7))
    tf = tagline_box.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(2)
    p = tf.paragraphs[0]
    p.font.size = Pt(13)
    run1 = p.add_run()
    run1.text = "读取已采集的 Ascend/CANN msprof 产物，将 profiler 事件整理成"
    run1.font.size = Pt(13)
    run1.font.color.rgb = TEXT_SECONDARY
    run1.font.name = "Arial"
    run2 = p.add_run()
    run2.text = "可读、可查、可比较"
    run2.font.size = Pt(13)
    run2.font.color.rgb = TEXT_PRIMARY
    run2.font.bold = True
    run2.font.name = "Arial"
    run3 = p.add_run()
    run3.text = "的性能地图。"
    run3.font.size = Pt(13)
    run3.font.color.rgb = TEXT_SECONDARY
    run3.font.name = "Arial"

    # --- Section: 核心能力 ---
    add_text_box(slide, left_x, Inches(2.15), Inches(2), Inches(0.35),
                 "核心能力", font_size=11, bold=True, color=TEXT_MUTED)

    # Feature list
    features = [
        (BLUE, "语义锚点树", "从 profiler 事件恢复循环结构和重复模式，生成可读 tree-map"),
        (GREEN, "增强 SQLite DB", "原始表不动，注入 traceloom_* 视图，用 SQL 无限深挖"),
        (PURPLE, "成本归因模型", "compute / communication / aux / idle 逐节点拆分成本"),
        (BLUE, "多卡对比分析", "自动发现设备，按 device 过滤或全局 top-N 热点分析"),
        (GREEN, "轻量离线设计", "纯分析工具，不安装驱动、不依赖 Docker、不绑定运行时"),
    ]

    y_start = Inches(2.5)
    for i, (icon_color, title, desc) in enumerate(features):
        y = y_start + Inches(i * 0.55)

        # Diamond icon
        diamond = slide.shapes.add_shape(
            MSO_SHAPE.DIAMOND, left_x, y + Pt(3), Pt(14), Pt(14)
        )
        diamond.fill.solid()
        diamond.fill.fore_color.rgb = icon_color
        diamond.line.fill.background()

        # Feature text
        feat_box = slide.shapes.add_textbox(Inches(0.95), y, Inches(5.5), Inches(0.5))
        ftf = feat_box.text_frame
        ftf.word_wrap = True
        ftf.margin_top = Pt(0)
        p = ftf.paragraphs[0]
        r1 = p.add_run()
        r1.text = title
        r1.font.size = Pt(12)
        r1.font.bold = True
        r1.font.color.rgb = TEXT_PRIMARY
        r1.font.name = "Arial"
        r2 = p.add_run()
        r2.text = f"  —  {desc}"
        r2.font.size = Pt(11)
        r2.font.color.rgb = TEXT_SECONDARY
        r2.font.name = "Arial"

    # --- Section: 面向用户 ---
    add_text_box(slide, left_x, Inches(5.35), Inches(2), Inches(0.3),
                 "面向用户", font_size=11, bold=True, color=TEXT_MUTED)

    users_box = slide.shapes.add_textbox(left_x, Inches(5.65), Inches(5.8), Inches(0.8))
    utf = users_box.text_frame
    utf.word_wrap = True
    utf.margin_top = Pt(0)
    users = [
        "做分布式推理、通信优化、算子优化的研发人员",
        "已有可运行的 Ascend 环境和 msprof 产物",
        "想快速定位热点循环、分析时间构成，并能用 SQL 继续深挖",
    ]
    for j, u in enumerate(users):
        if j == 0:
            p = utf.paragraphs[0]
        else:
            p = utf.add_paragraph()
        p.space_before = Pt(3)
        r = p.add_run()
        r.text = f"•  {u}"
        r.font.size = Pt(11)
        r.font.color.rgb = TEXT_SECONDARY
        r.font.name = "Arial"

    # ===================== RIGHT COLUMN =====================
    right_x = Inches(6.9)

    # --- Section: 推荐工作流 ---
    add_text_box(slide, right_x, Inches(1.3), Inches(3), Inches(0.3),
                 "推荐工作流", font_size=11, bold=True, color=TEXT_MUTED)

    # Workflow card background
    wf_card = add_rounded_rect(slide, right_x, Inches(1.6), Inches(5.7), Inches(2.1),
                               BG_CARD, BORDER)

    # Workflow steps
    steps = [
        ("1", "traceloom analyze /path/to/msprof_raw"),
        ("2", "阅读 summary.md + tree-map.md 定位全局视图"),
        ("3", "定位热点 node（如 N060，按 total_us 排序）"),
        ("4", "queries/*.sql 深挖事件、occurrence、成本构成"),
    ]

    for i, (num, text) in enumerate(steps):
        y = Inches(1.72) + Inches(i * 0.47)

        # Step number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, right_x + Inches(0.2), y + Pt(2), Pt(18), Pt(18)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(0x12, 0x24, 0x3A)
        circle.line.fill.background()
        circle.text_frame.paragraphs[0].text = num
        circle.text_frame.paragraphs[0].font.size = Pt(9)
        circle.text_frame.paragraphs[0].font.bold = True
        circle.text_frame.paragraphs[0].font.color.rgb = BLUE
        circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        circle.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Step text
        add_text_box(slide, right_x + Inches(0.55), y, Inches(5.0), Inches(0.4),
                     text, font_size=11, color=TEXT_SECONDARY,
                     font_name="Consolas")

        # Arrow between steps
        if i < len(steps) - 1:
            arrow_y = y + Inches(0.35)
            add_text_box(slide, right_x + Inches(0.22), arrow_y, Inches(0.3), Inches(0.25),
                         "↓", font_size=12, color=TEXT_MUTED)

    # --- Section: 输出产物 ---
    add_text_box(slide, right_x, Inches(3.9), Inches(3), Inches(0.3),
                 "输出产物", font_size=11, bold=True, color=TEXT_MUTED)

    # Output cards (2x2 grid)
    outputs = [
        ("tree-map.md", "性能地图 / 热点循环树"),
        ("*.augmented.db", "增强 SQLite / 可查询"),
        ("summary.md", "全局概览 / 设备信息"),
        ("queries/*.sql", "内置分析查询脚本"),
    ]

    card_w = Inches(2.7)
    card_h = Inches(0.7)
    gap_x = Inches(0.15)
    gap_y = Inches(0.1)

    for i, (name, desc) in enumerate(outputs):
        col = i % 2
        row = i // 2
        cx = right_x + col * (card_w + gap_x)
        cy = Inches(4.2) + row * (card_h + gap_y)

        card = add_rounded_rect(slide, cx, cy, card_w, card_h, BG_CARD, BORDER)
        ctf = card.text_frame
        ctf.word_wrap = True
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ctf.margin_left = Pt(10)

        p = ctf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = name
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = TEXT_PRIMARY
        r.font.name = "Consolas"

        p2 = ctf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(2)
        r2 = p2.add_run()
        r2.text = desc
        r2.font.size = Pt(9)
        r2.font.color.rgb = TEXT_MUTED
        r2.font.name = "Arial"

    # --- Section: TraceLoom 不做什么 ---
    add_text_box(slide, right_x, Inches(5.65), Inches(3), Inches(0.3),
                 "边界说明", font_size=11, bold=True, color=TEXT_MUTED)

    boundary_box = slide.shapes.add_textbox(right_x, Inches(5.92), Inches(5.7), Inches(0.9))
    btf = boundary_box.text_frame
    btf.word_wrap = True
    btf.margin_top = Pt(0)
    boundaries = [
        "不安装驱动、CANN、torch 或 vLLM — 纯离线分析",
        "不绑定特定运行时或容器环境",
        "不提交大型 profiler DB — 只产出轻量 CSV/MD/SQL",
    ]
    for j, b in enumerate(boundaries):
        if j == 0:
            p = btf.paragraphs[0]
        else:
            p = btf.add_paragraph()
        p.space_before = Pt(3)
        r = p.add_run()
        r.text = f"✗  {b}"
        r.font.size = Pt(10)
        r.font.color.rgb = TEXT_MUTED
        r.font.name = "Arial"

    # ===================== FOOTER =====================
    # Footer separator
    footer_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(6.95), Inches(11.9), Pt(1)
    )
    footer_line.fill.solid()
    footer_line.fill.fore_color.rgb = BORDER
    footer_line.line.fill.background()

    # Organization
    add_text_box(slide, Inches(0.7), Inches(7.05), Inches(5), Inches(0.35),
                 "vLLM-HUST  ·  分布式推理性能分析基础设施", font_size=11,
                 bold=False, color=TEXT_MUTED)

    # Tech tags
    tags = ["Ascend/CANN", "msprof", "SQLite", "Offline", "Python"]
    tag_x = Inches(9.2)
    for t in tags:
        tag_shape = add_rounded_rect(slide, tag_x, Inches(7.05), Inches(0.9), Inches(0.3),
                                     RGBColor(0x0E, 0x1F, 0x12), GREEN)
        tag_shape.text_frame.paragraphs[0].text = t
        tag_shape.text_frame.paragraphs[0].font.size = Pt(8)
        tag_shape.text_frame.paragraphs[0].font.color.rgb = GREEN
        tag_shape.text_frame.paragraphs[0].font.bold = True
        tag_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        tag_shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        tag_x += Inches(0.95)


def main():
    prs = Presentation()
    # Set widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_slide(prs)

    out_path = os.path.join(os.path.dirname(__file__), "TraceLoom-Introduction.pptx")
    prs.save(out_path)
    print(f"Saved: {out_path}")


if __name__ == "__main__":
    main()
